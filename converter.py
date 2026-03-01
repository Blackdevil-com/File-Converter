import subprocess
import os
import fitz
from pdf2docx import Converter
from PIL import Image
from pdf2image import convert_from_path
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import pdfkit
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from xlsx2html import xlsx2html


def convert_to_pdf(input_path, output_dir, timeout=30):
    os.makedirs(output_dir, exist_ok=True)

    command = [
        "soffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        input_path
    ]

    try:
        subprocess.run(command, check=True, timeout=timeout)
    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out")
    except subprocess.CalledProcessError:
        raise Exception("LibreOffice conversion failed")

    filename = os.path.basename(input_path)
    pdf_filename = os.path.splitext(filename)[0] + ".pdf"
    return os.path.join(output_dir, pdf_filename)


def convert_pdf_to_word(input_path, output_path):
    try:
        cv = Converter(input_path)
        cv.convert(output_path)
        cv.close()
    except Exception:
        raise Exception("PDF to Word conversion failed")

def convert_multiple_images_to_pdf(
    input_paths: list,
    output_path: str,
    orientation: str = "portrait",
    scale: str = "full"
):
    # A4 dimensions at 150 DPI
    A4_PORTRAIT = (1240, 1754)
    A4_LANDSCAPE = (1754, 1240)

    page_width, page_height = A4_PORTRAIT
    if orientation == "landscape":
        page_width, page_height = A4_LANDSCAPE

    scale_map = {"full": 1.0, "medium": 0.7, "small": 0.5}
    if scale not in scale_map:
        raise ValueError("Invalid scale: full, medium, small")
    img_scale = scale_map[scale]

    pdf_pages = []

    try:
        for path in input_paths:
            img = Image.open(path)

            # Convert to RGB
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")

            # Max image size inside page
            max_width = int(page_width * img_scale)
            max_height = int(page_height * img_scale)

            # Resize image proportionally
            img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)

            # Create blank A4 page
            page = Image.new("RGB", (page_width, page_height), "white")

            # Center image
            x = (page_width - img.width) // 2
            y = (page_height - img.height) // 2
            page.paste(img, (x, y))

            pdf_pages.append(page)

        # Save all pages to single PDF
        if pdf_pages:
            pdf_pages[0].save(
                output_path,
                save_all=True,
                append_images=pdf_pages[1:],
                resolution=150.0,
                quality=95
            )
        else:
            raise Exception("No images to convert")

    except Exception as e:
        raise Exception(f"Multi-image to PDF failed: {str(e)}")
    
def pdf_pages_to_images(input_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    try:
        images = convert_from_path(input_path, dpi=200)
        image_paths = []
        for i, img in enumerate(images):
            img_path = os.path.join(output_dir, f"page_{i+1}.jpg")
            img.save(img_path, "JPEG")
            image_paths.append(img_path)
        return image_paths
    except Exception as e:
        raise Exception(f"PDF to Page Images failed: {str(e)}")
    

# Mode 2: Extract embedded images from PDF
def pdf_extract_images(input_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    try:
        doc = fitz.open(input_path)
        img_count = 0
        image_paths = []

        for page_index in range(len(doc)):
            page = doc[page_index]
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                ext = base_image["ext"]
                img_count += 1
                img_path = os.path.join(output_dir, f"image_{img_count}.{ext}")
                with open(img_path, "wb") as f:
                    f.write(image_bytes)
                image_paths.append(img_path)

        return image_paths
    except Exception as e:
        raise Exception(f"PDF Extract Images failed: {str(e)}")
    

def create_text_watermark(text, output_path, page_size=letter):
    c = canvas.Canvas(output_path, pagesize=page_size)
    width, height = page_size
    c.setFont("Helvetica", 40)
    c.setFillAlpha(0.3)  # transparency
    c.saveState()
    c.translate(width/2, height/2)
    c.rotate(45)  # diagonal watermark
    c.drawCentredString(0, 0, text)
    c.restoreState()
    c.save()



def excel_to_pdf(input_excel, output_pdf):
    """
    Converts an Excel file to PDF.
    """
    temp_html = output_pdf.replace(".pdf", ".html")

    # Convert Excel to HTML
    xlsx2html(input_excel, temp_html)

    # Convert HTML to PDF
    pdfkit.from_file(temp_html, output_pdf)

    # Remove temporary HTML
    os.remove(temp_html)
    return output_pdf



def pptx_to_pdf(input_pptx, output_pdf):
    prs = Presentation(input_pptx)
    c = canvas.Canvas(output_pdf, pagesize=A4)
    width, height = A4

    for slide in prs.slides:
        # Save slide as image temporarily
        img_path = "temp_slide.png"
        slide.shapes._spTree.write(img_path)  # only works in some environments
        if os.path.exists(img_path):
            slide_img = Image.open(img_path)
            slide_img = slide_img.resize((int(width), int(height)))
            slide_img.save(img_path)
            c.drawImage(img_path, 0, 0, width=width, height=height)
            os.remove(img_path)
        c.showPage()

    c.save()
    return output_pdf